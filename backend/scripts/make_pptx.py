"""Generate the two KI-Challenge presentations as .pptx (python-pptx).

Run inside the backend container (has Python + pip):
    pip install python-pptx
    python /app/scripts/make_pptx.py /out

Produces:
    Investa_Offer_Intelligence_Technical.pptx
    Investa_Offer_Intelligence_Business.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

# --- Investa palette -------------------------------------------------------- #
PETROL = RGBColor(0x0E, 0x3B, 0x4C)
PETROL_DARK = RGBColor(0x08, 0x22, 0x2C)
ACCENT = RGBColor(0x2A, 0x70, 0x88)
SAND = RGBColor(0xC8, 0xA1, 0x5A)
INK = RGBColor(0x0B, 0x2D, 0x3A)
SLATE = RGBColor(0x55, 0x6A, 0x74)
LIGHT = RGBColor(0xEE, 0xF2, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x0E, 0x7A, 0x5A)
ROSE = RGBColor(0x9F, 0x44, 0x56)

# 16:9
SW = Emu(12192000)
SH = Emu(6858000)
MARGIN = Emu(685800)  # ~0.75"


def _slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE

    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _fill(sh, color)
    sh.shadow.inherit = False
    return sh


def _set(run, size, color, bold=False, italic=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def title_slide(prs, eyebrow, title, subtitle, footer):
    s = _slide(prs)
    _rect(s, 0, 0, SW, SH, PETROL)
    _rect(s, 0, Emu(5400000), SW, Emu(60000), SAND)

    _, tf = _box(s, MARGIN, Emu(1500000), Emu(10800000), Emu(600000))
    r = tf.paragraphs[0].add_run()
    r.text = eyebrow
    _set(r, 16, SAND, bold=True)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    _, tf = _box(s, MARGIN, Emu(2150000), Emu(10800000), Emu(2000000))
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 40, WHITE, bold=True, font="Georgia")

    _, tf = _box(s, MARGIN, Emu(4000000), Emu(10800000), Emu(900000))
    r = tf.paragraphs[0].add_run()
    r.text = subtitle
    _set(r, 18, RGBColor(0x8F, 0xB8, 0xC6))

    _, tf = _box(s, MARGIN, Emu(6050000), Emu(10800000), Emu(500000))
    r = tf.paragraphs[0].add_run()
    r.text = footer
    _set(r, 12, RGBColor(0x8F, 0xB8, 0xC6))


def content_slide(prs, kicker, title, bullets, note=None):
    """bullets: list of (text, level, bold) or (text, level) or str."""
    s = _slide(prs)
    _rect(s, 0, 0, SW, Emu(1150000), PETROL)
    _rect(s, 0, Emu(1150000), SW, Emu(28000), SAND)

    _, tf = _box(s, MARGIN, Emu(210000), Emu(11000000), Emu(360000))
    r = tf.paragraphs[0].add_run()
    r.text = kicker
    _set(r, 13, SAND, bold=True)

    _, tf = _box(s, MARGIN, Emu(560000), Emu(11000000), Emu(560000))
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 26, WHITE, bold=True, font="Georgia")

    top = Emu(1450000)
    _, tf = _box(s, MARGIN, top, Emu(10820000), Emu(4900000))
    tf.word_wrap = True
    first = True
    for b in bullets:
        if isinstance(b, str):
            text, level, bold = b, 0, False
        elif len(b) == 2:
            text, level = b
            bold = False
        else:
            text, level, bold = b
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(6)
        run = p.add_run()
        bullet_char = "▪ " if level == 0 else "– "
        run.text = bullet_char + text
        _set(run, 19 if level == 0 else 16, INK if level == 0 else SLATE, bold=bold)

    if note:
        nb = _rect(s, MARGIN, Emu(5950000), Emu(10820000), Emu(620000), LIGHT)
        nb.line.color.rgb = ACCENT
        nb.line.width = Pt(1)
        tf = nb.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(10)
        tf.margin_top = Pt(6)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "→ " + note
        _set(r, 13, PETROL, italic=True)


def table_slide(prs, kicker, title, headers, rows, note=None):
    s = _slide(prs)
    _rect(s, 0, 0, SW, Emu(1150000), PETROL)
    _rect(s, 0, Emu(1150000), SW, Emu(28000), SAND)

    _, tf = _box(s, MARGIN, Emu(210000), Emu(11000000), Emu(360000))
    r = tf.paragraphs[0].add_run()
    r.text = kicker
    _set(r, 13, SAND, bold=True)
    _, tf = _box(s, MARGIN, Emu(560000), Emu(11000000), Emu(560000))
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 26, WHITE, bold=True, font="Georgia")

    ncols = len(headers)
    nrows = len(rows) + 1
    left, top = MARGIN, Emu(1500000)
    width = Emu(10820000)
    height = Emu(4200000)
    gtbl = s.shapes.add_table(nrows, ncols, left, top, width, height).table

    for c, h in enumerate(headers):
        cell = gtbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PETROL
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        _set(run, 14, WHITE, bold=True)
    for r_i, row in enumerate(rows, start=1):
        for c_i, val in enumerate(row):
            cell = gtbl.cell(r_i, c_i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_i % 2 else LIGHT
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            _set(run, 13, INK)

    if note:
        _, tf = _box(s, MARGIN, Emu(5950000), Emu(10820000), Emu(600000))
        r = tf.paragraphs[0].add_run()
        r.text = "→ " + note
        _set(r, 13, PETROL, italic=True)


def section_slide(prs, number, title, subtitle=None):
    s = _slide(prs)
    _rect(s, 0, 0, SW, SH, PETROL_DARK)
    _rect(s, MARGIN, Emu(2900000), Emu(900000), Emu(120000), SAND)
    _, tf = _box(s, MARGIN, Emu(2150000), Emu(10000000), Emu(700000))
    r = tf.paragraphs[0].add_run()
    r.text = number
    _set(r, 16, SAND, bold=True)
    _, tf = _box(s, MARGIN, Emu(3150000), Emu(10800000), Emu(1200000))
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 34, WHITE, bold=True, font="Georgia")
    if subtitle:
        _, tf = _box(s, MARGIN, Emu(4350000), Emu(10800000), Emu(800000))
        r = tf.paragraphs[0].add_run()
        r.text = subtitle
        _set(r, 18, RGBColor(0x8F, 0xB8, 0xC6))


def _img_size(path: Path):
    """Return (w, h) in px from a PNG header — no PIL dependency."""
    import struct

    with open(path, "rb") as f:
        head = f.read(24)
    if head[:8] == b"\x89PNG\r\n\x1a\n":
        w, h = struct.unpack(">II", head[16:24])
        return w, h
    return 1440, 900


def screenshot_slide(prs, kicker, title, img_path, callouts=None):
    """Slide with a header + a framed screenshot, fit to area, plus optional callouts."""
    img_path = Path(img_path)
    s = _slide(prs)
    _rect(s, 0, 0, SW, Emu(1150000), PETROL)
    _rect(s, 0, Emu(1150000), SW, Emu(28000), SAND)
    _rect(s, 0, Emu(1178000), SW, Emu(5680000), LIGHT)

    _, tf = _box(s, MARGIN, Emu(210000), Emu(11000000), Emu(360000))
    r = tf.paragraphs[0].add_run()
    r.text = kicker
    _set(r, 13, SAND, bold=True)
    _, tf = _box(s, MARGIN, Emu(560000), Emu(11000000), Emu(560000))
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 26, WHITE, bold=True, font="Georgia")

    if not img_path.exists():
        return

    w_px, h_px = _img_size(img_path)
    aspect = w_px / h_px

    has_callouts = bool(callouts)
    area_left = MARGIN
    area_top = Emu(1380000)
    area_w = Emu(7600000 if has_callouts else 10820000)
    area_h = Emu(5200000)

    # Fit by the tighter dimension.
    max_w = int(area_w)
    max_h = int(area_h)
    if max_w / aspect <= max_h:
        draw_w = max_w
        draw_h = int(max_w / aspect)
    else:
        draw_h = max_h
        draw_w = int(max_h * aspect)
    # cap very tall images
    if draw_h > max_h:
        draw_h = max_h
        draw_w = int(max_h * aspect)

    left = int(area_left) + (int(area_w) - draw_w) // 2
    top = int(area_top) + (int(area_h) - draw_h) // 2

    frame = _rect(s, Emu(left - 18000), Emu(top - 18000), Emu(draw_w + 36000), Emu(draw_h + 36000), WHITE)
    frame.line.color.rgb = RGBColor(0xCB, 0xD5, 0xDB)
    frame.line.width = Pt(1)
    s.shapes.add_picture(str(img_path), Emu(left), Emu(top), Emu(draw_w), Emu(draw_h))

    if has_callouts:
        _, tf = _box(s, Emu(8400000), Emu(1500000), Emu(3300000), Emu(4800000))
        tf.word_wrap = True
        first = True
        for c in callouts:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(10)
            run = p.add_run()
            run.text = "● " + c
            _set(run, 15, INK)


# ============================================================================ #
# TECHNICAL DECK
# ============================================================================ #
def build_technical(out: Path, shots: Path):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    title_slide(
        prs,
        "KI CHALLENGE · FORWARD DEPLOYED ENGINEER (AI / GenAI)",
        "Investa Offer Intelligence",
        "Technische Präsentation — Architektur, Umsetzung, Entscheidungslogik & Trade-offs",
        "Bassem Trifa   ·   Panel: Herzhoff (CIO), Hesse, Quintes (anyon.studio), Vojinovic",
    )

    content_slide(
        prs, "AGENDA", "Was Sie gleich sehen",
        [
            ("Live-Demo zuerst — die ganze Lösung an einer echten Makler-E-Mail", 0, True),
            ("Architektur: 3 Services, ein Befehl (docker compose up)", 0),
            ("Pipeline & eingesetzte Technologien / LLMs", 0),
            ("Prompt Engineering — nachvollziehbar & halluzinationsarm", 0),
            ("Entscheidungslogik: Dedup + asset-klassen-bewusstes Scoring", 0),
            ("Trade-offs, Performance, Skalierung & KI im Entwicklungsprozess", 0),
        ],
        note="Show, then tell — erst die Demo, dann die Technik dahinter.",
    )

    section_slide(prs, "01", "Live-Demo", "Eine echte .msg-Angebots-E-Mail, Ende zu Ende")

    screenshot_slide(
        prs, "DEMO · DASHBOARD", "Portfolio-Überblick: Posteingang als Rangliste",
        shots / "home_hero.png",
        ["Kennzahlen auf einen Blick", "Score-Verteilung", "Top-Angebote priorisiert"],
    )
    screenshot_slide(
        prs, "DEMO · ANGEBOT", "Ein Angebot, Ende zu Ende",
        shots / "detail_hero.png",
        ["Extrahierte Fakten", "Karte + Marktdaten", "Asset-Klasse + Score"],
    )

    content_slide(
        prs, "PROBLEM", "Worum es geht",
        [
            ("Investa erhält laufend Maklerangebote (E-Mail + PDF) — ein realer Deal-Kanal", 0),
            ("Aber: stark heterogen — viele sind bereits bekannt, unwirtschaftlich oder unvollständig", 0),
            ("Manuelle Prüfung ist langsam, teuer und bindet Fachressourcen", 0),
            ("Ziel: automatisierte Triage → extrahieren, Dubletten erkennen, anreichern, 0–10 bewerten", 0, True),
            ("Mein Framing: nicht „ein Angebot bewerten“, sondern „den ganzen Posteingang triagieren“", 0),
        ],
    )

    content_slide(
        prs, "ARCHITEKTUR", "Drei Services, ein Befehl",
        [
            ("React (nginx)  →  FastAPI  →  PostgreSQL (pg_trgm / PostGIS)", 0, True),
            ("Pipeline = asynchroner Background-Task mit persistiertem Status (UI pollt)", 0),
            ("received → parsing → extracting → matching → enriching → scoring → done|failed", 1),
            ("Provider-agnostischer LLM-Client (OpenAI-kompatibel)", 0, True),
            ("Default: GitHub Models (kostenlos); Wechsel zu OpenAI/Azure/Gemini/Groq = 3 Zeilen .env", 1),
            ("Docker Compose: Healthchecks, Volumes für Persistenz, lokal lauffähig", 0),
        ],
        note="Kein Vendor-Lock-in, kein Lizenzproblem — Provider-Wechsel ist reine Konfiguration.",
    )

    table_slide(
        prs, "PIPELINE", "Stufen & eingesetzte Technologien",
        ["Stufe", "Aufgabe", "Technologie"],
        [
            ["Ingestion", ".msg/.eml/PDF/Bild → Text (+OCR)", "extract-msg, pdfplumber, pypdfium2"],
            ["Extraction", "Text → strukturierte JSON-Fakten", "LLM, structured output, Pydantic"],
            ["Dedup", "Bekannt-Objekt-Prüfung", "pg_trgm + RapidFuzz + Geo"],
            ["Enrichment", "Geo, POIs, Demografie, Miete", "Nominatim, Overpass, Wikidata, Seed"],
            ["Bildanalyse", "gerenderte Seiten → Hinweise", "Vision-LLM"],
            ["Scoring", "asset-klassen-bewusst 0–10", "Heuristik + LLM"],
            ["Reply", "„bereits bekannt“-Entwurf", "LLM"],
        ],
    )

    content_slide(
        prs, "PROMPT ENGINEERING", "Nachvollziehbar & halluzinationsarm",
        [
            ("Prompts sind versionierte Dateien in app/llm/prompts/ — keine Inline-Strings", 0),
            ("Drei fokussierte Prompts (Extraktion / Scoring / Reply) statt eines Mega-Prompts", 0),
            ("Anti-Halluzination: striktes JSON-Schema, null für Unbekanntes + confidence, T=0.1", 0, True),
            ("LLM ≠ Taschenrechner: Code rechnet (Rendite, €/m², Distanzen) — LLM urteilt", 0, True),
            ("Volle Traceability: jeder Call speichert {system,user,model,params,raw_response}", 0),
            ("→ im UI-Panel „Prompt-Trace“ sichtbar: jede Zahl rückführbar auf den exakten Prompt", 1),
        ],
        note="Kernfrage Cedric: Wie verhinderst du Halluzinationen? Schema + null + confidence + low temp.",
    )

    content_slide(
        prs, "ENTSCHEIDUNGSLOGIK 1", "Dedup — hybrid & erklärbar",
        [
            ("Adressen sind unsauber (Bürgerstr. vs. Bürgerstraße 44) → kein Einzelverfahren reicht", 0),
            ("normalize → pg_trgm Trigramm-Vorfilter (schnell, in SQL)", 0),
            ("→ RapidFuzz token_sort_ratio (Feinabgleich) → Geo-Haversine-Gate", 1),
            ("Match wenn addr_sim ≥ 0.88  ODER  (≥ 0.75 und < 150 m); Graubereich → needs_review", 0, True),
            ("Schwellen konfigurierbar; Match-Evidenz gespeichert", 0),
            ("Demo: Re-Upload Bürgerstraße/Bernau → „bereits bekannt“ → Makler-Antwort, keine Provision", 0),
        ],
    )

    section_slide(
        prs, "02", "Asset-klassen-bewusstes Scoring",
        "Der Differenzierer — kein Einheits-Kriterienset",
    )

    content_slide(
        prs, "SCORING — KERNIDEE", "Jede Asset-Klasse mit den richtigen Kriterien",
        [
            ("Der Fehler im naiven Design: ein fixes Kriterienset für alles ist fachlich falsch", 0),
            ("Ein Baugrundstück, ein vermietetes MFH und ein Büro werden anders bewertet", 0, True),
            ("Lösung: klassifizieren → Profil wählt, WELCHE Sub-Scores gelten und ihre Gewichte", 0, True),
            ("Hybrid: Heuristik-Sub-Scores + LLM-Urteil (Lage); Risiko-Abzug für strukturelle Mängel", 0),
            ("Top-3-Treiber sichtbar; Profil + Gewichte in scoring.yaml — ohne Code tunebar", 0),
        ],
        note="Adressiert direkt „Relevanz der gewählten Kriterien“ und Investas Multi-Asset-Realität.",
    )

    table_slide(
        prs, "SCORING — LIVE-BELEG", "Dieselben Angebote, unterschiedliche Kriterien",
        ["Angebot", "Erkannt als", "Score", "Bewertet auf"],
        [
            ["Raintal-Höfe (München)", "Entwicklungsgrundstück", "6.45", "Bebaubarkeit (111 WE), Lage  · −1.2 Genehmigungsrisiko"],
            ["Geesthacht (24 WE)", "Einkommensobjekt", "6.96", "Rendite, Vermietung 9, Reversion 2.0"],
            ["Bad Homburg (Büro)", "Gewerbeobjekt", "8.85", "Rendite, Vermietung, Lage, Zustand"],
        ],
        note="Raintal wird NICHT mehr auf Rendite bewertet; Geesthacht-Reversion niedrig, weil Miete bereits über Benchmark — echte Nuance.",
    )

    screenshot_slide(
        prs, "SCORING · UI", "Transparenz im Detail — Sub-Scores, Treiber, Risikoabzug",
        shots / "detail.png",
        ["„Bewertet als …“-Badge", "Sub-Scores + Gewichte", "Top-Treiber", "Risikoabzug", "Prompt-Trace"],
    )

    content_slide(
        prs, "DATEN & RESILIENZ", "Externe Anreicherung (kostenlos, gecacht)",
        [
            ("Nominatim (Geocoding, löst auch unvollständige Adressen)", 0),
            ("Overpass (POIs/ÖPNV), Wikidata (Demografie), Mietspiegel-Seed (dokumentierte Vereinfachung)", 0),
            ("Jeder externe Call in DB gecacht mit source + fetched_at", 0, True),
            ("Pro-Host-Throttling + Retry/Backoff; aussagekräftiger User-Agent", 0),
            ("War-Story: OSM blockt Platzhalter-User-Agents (403) — gefunden & gefixt", 1),
        ],
    )

    table_slide(
        prs, "PERFORMANCE", "Gemessen auf diesem Build",
        ["Pfad", "Zeit / Footprint"],
        [
            ["Volle Pipeline (3 LLM-Calls + 4 APIs + 4-Bild-Vision)", "~22 s"],
            ["Bekanntes Objekt (Dublette, Short-Circuit)", "~9 s"],
            ["Footprint im Leerlauf", "Backend 487 MB · DB 87 MB · Frontend 9 MB"],
        ],
        note="Vision (~11 s) ist größter Einzelposten und optional (vision_enabled); next step: parallel zur Anreicherung.",
    )

    table_slide(
        prs, "TRADE-OFFS", "Bewusste Entscheidungen → Produktions-Schritt",
        ["Entscheidung", "Warum (pragmatisch)", "Nächster Schritt"],
        [
            ["Async-Task statt Queue", "weniger Teile; Status persistiert", "Celery/Redis + Worker"],
            ["Mietspiegel-Seed-CSV", "keine freie Voll-API", "lizenzierter Markt-Feed"],
            ["Keyword-Klassifikator", "deterministisch, transparent", "LLM-Klassifikator + Confidence"],
            ["create_all statt Migrationen", "ok für 3-Tage-Prototyp", "Alembic"],
            ["Antwort nur Entwurf", "laut Aufgabe out of scope", "Graph/SMTP + Freigabe"],
        ],
        note="Trade-offs als bewusste Entscheidungen mit Ausblick — nicht als Lücken.",
    )

    content_slide(
        prs, "KI IM PROZESS & SKALIERUNG", "Wie es gebaut wurde — und wie es wächst",
        [
            ("Gebaut mit KI-Coding-Agent (GitHub Copilot) im Pair-Programming-Loop", 0),
            ("KI beschleunigte Boilerplate, Parsing-Edge-Cases, Refactors", 1),
            ("Menschliches Urteil trieb Scoring-Design, Asset-Klassen-Modell, Trade-offs", 1, True),
            ("Skalierung: stateless API → horizontal; Pipeline → Queue + Worker", 0),
            ("Externe Calls gecacht; LLM provider-swappbar; API-first für SAP/SharePoint/Mietspiegel", 0, True),
        ],
    )

    content_slide(
        prs, "ABSCHLUSS", "Zusammenfassung",
        [
            ("Ende-zu-Ende, automatisiert, nachvollziehbar, asset-klassen-bewusst", 0, True),
            ("In 3 Tagen gebaut, läuft mit einem Befehl", 0),
            ("Designed, um in Investas Datenlandschaft einzudocken", 0),
            ("Vielen Dank — Fragen & Austausch", 0, True),
        ],
    )

    prs.save(str(out))
    print("wrote", out)


# ============================================================================ #
# BUSINESS DECK
# ============================================================================ #
def build_business(out: Path, shots: Path):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    title_slide(
        prs,
        "KI CHALLENGE · FORWARD DEPLOYED ENGINEER (AI / GenAI)",
        "Investa Offer Intelligence",
        "Fachliche Präsentation — KI-gestützte Triage & Bewertung von Maklerangeboten",
        "Bassem Trifa   ·   Panel: Femerling (Geschäftsführer), Vojinovic   ·   ~30 Min",
    )

    content_slide(
        prs, "DAS PROBLEM", "In Investas Worten",
        [
            ("Investa erhält laufend Maklerangebote per E-Mail + PDF — ein realer Deal-Kanal", 0),
            ("Aber stark heterogen:", 0),
            ("viele bereits bekannt (Re-Angebote → Provisionsstreit-Risiko)", 1),
            ("viele wirtschaftlich unattraktiv", 1),
            ("viele unvollständig oder schwer vergleichbar", 1),
            ("Heute liest eine Fachkraft jedes Angebot manuell — langsam, teuer, inkonsistent", 0, True),
        ],
        note="Welche 3 Angebote dieser Woche sind die Zeit eines Analysten wert? Heute: alle lesen. Meine Lösung: in Sekunden.",
    )

    content_slide(
        prs, "DIE LÖSUNG", "Was ich gebaut habe",
        [
            ("Ein Tool, das jedes eingehende Angebot liest, prüft ob bereits bekannt,", 0, True),
            ("mit Marktdaten anreichert und 0–10 bewertet — mit klarer Empfehlung", 0, True),
            ("Aus einem vollen Posteingang wird eine kurze, priorisierte Liste", 0),
            ("Statt 30 PDFs zu lesen, schauen Sie auf die 3 grünen", 0, True),
        ],
    )

    section_slide(prs, "", "Live-Demo", "Ein echtes Angebot, Ende zu Ende")

    screenshot_slide(
        prs, "DEMO · ÜBERSICHT", "Aus vollem Posteingang wird eine Rangliste",
        shots / "home_hero.png",
        ["Kennzahlen", "Score-Verteilung", "Top-Angebote"],
    )
    screenshot_slide(
        prs, "DEMO · ANGEBOT", "Fakten, Karte, Bewertung — auf einen Blick",
        shots / "detail_hero.png",
        ["Automatisch extrahiert", "Mit Marktdaten angereichert", "Klare Empfehlung"],
    )

    table_slide(
        prs, "ERGEBNISSE", "Auf den echten Angeboten",
        ["Angebot", "Erkannt als", "Score", "Empfehlung"],
        [
            ["Geesthacht — vermietete Wohnanlage", "Einkommensobjekt", "6.96", "Prüfen"],
            ["Bad Homburg — Büro", "Gewerbeobjekt", "8.85", "Verfolgen"],
            ["Raintal-Höfe — Bauland (München)", "Entwicklungsgrundstück", "6.45", "Prüfen"],
            ["Hamburg — Grundstück m. Bestand", "Entwicklungsgrundstück", "4.70", "Prüfen"],
            ["Bernau / Bürgerstraße (erneut)", "Bereits bekannt", "—", "Auto-Antwort, keine Provision"],
        ],
        note="Der Punkt sind nicht die exakten Zahlen — jedes Angebot wurde als das verstanden, was es ist, und konsistent eingeordnet.",
    )

    content_slide(
        prs, "VERTRAUEN", "Warum die Bewertung verlässlich ist",
        [
            ("Ein Profi bewertet ein Baugrundstück anders als ein vermietetes MFH oder ein Büro", 0),
            ("Ein Tool, das alles gleich bewertet, wäre für Sie sichtbar falsch", 0, True),
            ("Daher erkennt das Tool zuerst den Objekttyp und wendet die richtigen Kriterien an:", 0),
            ("Bauland → Bebaubarkeit, Lage, Preis je künftiger Einheit", 1),
            ("Einkommensobjekt → Rendite, wie voll vermietet, Mietsteigerungspotenzial", 1),
            ("Es bestraft echte Risiken (Altlast, Erbpacht, Leerstand) und zeigt die Top-Treiber", 0, True),
        ],
        note="Es denkt wie ein Analyst — „ist es schon vermietet“ zählt beim MFH, ist beim Grundstück irrelevant.",
    )

    content_slide(
        prs, "KI IM EINSATZ", "Wo KI hilft — und wo nicht",
        [
            ("KI (das Sprachmodell) übernimmt das Urteil:", 0, True),
            ("messy E-Mails/PDFs lesen, Fakten ziehen, Lage/Risiken/Chancen einschätzen, Bilder lesen, Antworten entwerfen", 1),
            ("Klassische Berechnung übernimmt die Zahlen: Rendite, €/m², Distanzen", 0, True),
            ("Bewusst lasse ich die KI nicht rechnen — sie interpretiert, das System rechnet", 1),
            ("Für Sie: jede Zahl ist reproduzierbar und erklärbar — KI dort, wo sie Mehrwert schafft", 0),
        ],
    )

    content_slide(
        prs, "EHRLICHKEIT", "Grenzen & abgewogene Alternativen",
        [
            ("Entscheidungs-Unterstützung, kein Auto-Kauf — ein Mensch entscheidet immer", 0, True),
            ("Manche Daten sind öffentlich/approximativ (Referenzmiete = kuratierter Datensatz)", 0),
            ("Es kann irren — daher zeigt es Konfidenz, markiert dünne Datenlage und erklärt sich", 0),
            ("Abgewogen: reines Regelwerk (zu starr), reine KI-Bewertung (nicht auditierbar)", 0),
            ("Gewählt: der hybride Weg — klug UND erklärbar", 0, True),
        ],
    )

    content_slide(
        prs, "MEHRWERT & AUSBLICK", "Geschäftlicher Nutzen",
        [
            ("Schnellere Triage, keine guten Angebote übersehen, weniger Provisionsstreit", 0),
            ("Konsistente, dokumentierte Entscheidungen — Fachkräfte nur noch auf das Relevante", 0, True),
            ("Passt zu Investa: andockbar an SAP/SharePoint + echten Mietspiegel-Feed", 0),
            ("Nächste Schritte: interne Daten integrieren, „needs review“-Queue, aus Analysten-Feedback lernen", 0),
            ("In 3 Tagen ein lauffähiges Tool, das weiß, wo man zuerst hinschauen sollte", 0, True),
        ],
    )

    prs.save(str(out))
    print("wrote", out)


if __name__ == "__main__":
    out_dir = Path(sys.argv[1] if len(sys.argv) > 1 else ".")
    shots_dir = Path(sys.argv[2] if len(sys.argv) > 2 else (out_dir / "shots"))
    out_dir.mkdir(parents=True, exist_ok=True)
    build_technical(out_dir / "Investa_Offer_Intelligence_Technical.pptx", shots_dir)
    build_business(out_dir / "Investa_Offer_Intelligence_Business.pptx", shots_dir)
